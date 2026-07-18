"""
Core data processing engine built on Polars + PyArrow + DuckDB.
Handles reading, writing, and previewing all supported formats.
"""

import logging
import io
import os
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import polars as pl
import duckdb

logger = logging.getLogger(__name__)

# ─── Format groups ────────────────────────────────────────────────────────────

EXCEL_FORMATS  = {"xlsx", "xls", "xlsm", "xlsb"}
CSV_LIKE       = {"csv", "tsv", "txt"}
ARROW_FORMATS  = {"parquet", "feather"}
STRUCTURED     = {"json", "xml", "yaml", "yml"}
DB_FORMATS     = {"sqlite", "db"}
DOC_FORMATS    = {"docx", "pdf"}
ODS_FORMAT     = {"ods"}
PPTX_FORMATS   = {"pptx"}
HTML_FORMATS   = {"html", "htm"}
IMAGE_FORMATS  = {"jpg", "jpeg", "png", "bmp", "gif", "webp"}
SVG_FORMAT     = {"svg"}

# Conversions that bypass the DataFrame (direct binary transformation)
DIRECT_PAIRS: set[tuple[str, str]] = (
    {(img, "pdf")  for img in IMAGE_FORMATS} |
    {("pdf", img)  for img in ("jpg", "jpeg", "png")} |
    {("svg", "pdf"), ("pdf", "svg")}
)


class DataEngine:
    """
    Unified data engine using Polars as the primary processing backbone.
    Falls back to pandas / openpyxl where Polars lacks native support.
    """

    # ─── Sheet helpers ────────────────────────────────────────────────────────

    def get_excel_sheets(self, path: str, fmt: str) -> list[str]:
        """Return the list of worksheet names in an Excel file."""
        try:
            from python_calamine import CalamineWorkbook
            return CalamineWorkbook.from_path(path).sheet_names
        except Exception:
            pass
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            names = wb.sheetnames
            wb.close()
            return names
        except Exception:
            pass
        try:
            import pandas as pd
            return pd.ExcelFile(path).sheet_names
        except Exception:
            return []

    def read_all_sheets(
        self, path: str, fmt: str, sheets: list[str] | None = None
    ) -> dict[str, pl.DataFrame]:
        """Read selected (or all) sheets from an Excel file into a dict."""
        available = self.get_excel_sheets(path, fmt)
        to_read   = [s for s in sheets if s in available] if sheets else available
        return {name: self._read_excel(path, fmt, sheet=name) for name in to_read}

    def read_pdf_pages(self, path: str) -> dict[str, pl.DataFrame]:
        """Extract each PDF page as its own DataFrame (table or plain text)."""
        import pdfplumber
        result: dict[str, pl.DataFrame] = {}
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                label = f"صفحة {i}"
                rows: list[dict] = []
                headers: list[str] = []
                for tbl in page.extract_tables():
                    if not tbl:
                        continue
                    if not headers:
                        headers = [str(c or f"col_{j}") for j, c in enumerate(tbl[0])]
                    for row in tbl[1:]:
                        rows.append({h: str(c or "") for h, c in zip(headers, row)})
                if rows:
                    result[label] = pl.DataFrame(rows)
                else:
                    text  = page.extract_text() or ""
                    lines = [ln for ln in text.split("\n") if ln.strip()]
                    if lines:
                        result[label] = pl.DataFrame({"النص": lines})
        return result

    # ─── Read ─────────────────────────────────────────────────────────────────

    def read(self, path: str, fmt: str, sheet: str | int | None = None) -> pl.DataFrame:
        """Read any supported format into a Polars DataFrame."""
        fmt = fmt.lower().lstrip(".")

        if fmt in CSV_LIKE:
            sep = "\t" if fmt == "tsv" else ","
            return pl.read_csv(path, separator=sep, infer_schema_length=5000, ignore_errors=True)

        if fmt in EXCEL_FORMATS:
            return self._read_excel(path, fmt, sheet)

        if fmt == "parquet":
            return pl.read_parquet(path)

        if fmt == "feather":
            return pl.read_ipc(path)

        if fmt == "json":
            try:
                return pl.read_json(path)
            except Exception:
                return pl.read_ndjson(path)

        if fmt in {"yaml", "yml"}:
            return self._read_yaml(path)

        if fmt == "xml":
            return self._read_xml(path)

        if fmt in DB_FORMATS:
            return self._read_sqlite(path)

        if fmt == "docx":
            return self._read_docx_tables(path)

        if fmt == "pdf":
            return self._read_pdf_tables(path)

        if fmt == "ods":
            return self._read_ods(path)

        if fmt in PPTX_FORMATS:
            return self._read_pptx(path)

        if fmt in HTML_FORMATS:
            return self._read_html(path)

        if fmt in IMAGE_FORMATS:
            return self._read_image_meta(path)

        raise ValueError(f"Unsupported read format: {fmt}")

    # ── tabular readers ───────────────────────────────────────────────────────

    def _read_excel(self, path: str, fmt: str, sheet=None) -> pl.DataFrame:
        try:
            import python_calamine
            return pl.read_excel(path, sheet_name=sheet or 0, engine="calamine")
        except Exception:
            pass
        try:
            return pl.read_excel(path, sheet_name=sheet or 0, engine="openpyxl")
        except Exception:
            import pandas as pd
            if fmt == "xlsb":
                df_pd = pd.read_excel(path, sheet_name=sheet or 0, engine="pyxlsb")
            elif fmt == "xls":
                df_pd = pd.read_excel(path, sheet_name=sheet or 0, engine="xlrd")
            else:
                df_pd = pd.read_excel(path, sheet_name=sheet or 0)
            return pl.from_pandas(df_pd)

    def _read_yaml(self, path: str) -> pl.DataFrame:
        import yaml
        with open(path) as f:
            data = yaml.safe_load(f)
        if isinstance(data, list):
            return pl.DataFrame(data)
        return pl.DataFrame([data])

    def _read_xml(self, path: str) -> pl.DataFrame:
        import pandas as pd
        df_pd = pd.read_xml(path)
        return pl.from_pandas(df_pd)

    def _read_sqlite(self, path: str) -> pl.DataFrame:
        con = duckdb.connect()
        con.execute(f"ATTACH '{path}' AS src (TYPE sqlite)")
        tables = con.execute("SHOW TABLES").fetchall()
        if not tables:
            return pl.DataFrame()
        table_name = tables[0][0]
        return con.execute(f"SELECT * FROM src.{table_name}").pl()

    def _read_docx_tables(self, path: str) -> pl.DataFrame:
        from docx import Document
        doc = Document(path)
        rows: list[dict] = []
        headers: list[str] = []
        for table in doc.tables:
            if not headers:
                headers = [cell.text.strip() for cell in table.rows[0].cells]
            for row in table.rows[1:]:
                rows.append({h: c.text.strip() for h, c in zip(headers, row.cells)})
        return pl.DataFrame(rows) if rows else pl.DataFrame()

    def _read_pdf_tables(self, path: str) -> pl.DataFrame:
        import pdfplumber
        rows: list[dict] = []
        headers: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables():
                    if not table:
                        continue
                    if not headers:
                        headers = [str(c or f"col_{i}") for i, c in enumerate(table[0])]
                    for row in table[1:]:
                        rows.append({h: str(c or "") for h, c in zip(headers, row)})
        return pl.DataFrame(rows) if rows else pl.DataFrame()

    def _read_ods(self, path: str) -> pl.DataFrame:
        import pandas as pd
        df_pd = pd.read_excel(path, engine="odf")
        return pl.from_pandas(df_pd)

    def _read_pptx(self, path: str) -> pl.DataFrame:
        """Extract text content from each PPTX slide as a DataFrame row."""
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(path)
        rows: list[dict] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            # Also extract tables from slides
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    headers = [tbl.cell(0, c).text.strip() for c in range(tbl._tbl.col_count)]
                    for r in range(1, tbl._tbl.tr_count):
                        row_dict = {"slide": i}
                        for c, h in enumerate(headers):
                            row_dict[h or f"col_{c}"] = tbl.cell(r, c).text.strip()
                        rows.append(row_dict)
            if not any(r.get("slide") == i for r in rows):
                rows.append({"slide": i, "content": " | ".join(texts)})
        return pl.DataFrame(rows) if rows else pl.DataFrame()

    def _read_html(self, path: str) -> pl.DataFrame:
        """Read the first table found in an HTML file."""
        import pandas as pd
        tables = pd.read_html(path)
        if not tables:
            return pl.DataFrame()
        return pl.from_pandas(tables[0])

    def _read_image_meta(self, path: str) -> pl.DataFrame:
        """Return basic image metadata as a single-row DataFrame."""
        from PIL import Image
        img = Image.open(path)
        return pl.DataFrame([{
            "filename": Path(path).name,
            "format":   img.format or Path(path).suffix.lstrip(".").upper(),
            "mode":     img.mode,
            "width":    img.size[0],
            "height":   img.size[1],
        }])

    # ─── Write ────────────────────────────────────────────────────────────────

    def write(self, df: pl.DataFrame, path: str, fmt: str, **kwargs) -> None:
        """Write a Polars DataFrame to any supported output format."""
        fmt = fmt.lower().lstrip(".")

        if fmt == "csv":
            df.write_csv(path)
        elif fmt == "tsv":
            df.write_csv(path, separator="\t")
        elif fmt in EXCEL_FORMATS or fmt == "xlsx":
            df.write_excel(path)
        elif fmt == "parquet":
            df.write_parquet(path)
        elif fmt == "feather":
            df.write_ipc(path)
        elif fmt == "json":
            df.write_json(path)
        elif fmt == "ndjson":
            df.write_ndjson(path)
        elif fmt in HTML_FORMATS:
            df.to_pandas().to_html(path, index=False)
        elif fmt == "xml":
            df.to_pandas().to_xml(path, index=False)
        elif fmt in {"yaml", "yml"}:
            import yaml
            with open(path, "w") as f:
                yaml.dump(df.to_dicts(), f, allow_unicode=True)
        elif fmt == "ods":
            df.to_pandas().to_excel(path, engine="odf", index=False)
        elif fmt == "sqlite":
            con = duckdb.connect()
            con.register("df_table", df.to_arrow())
            con.execute(f"ATTACH '{path}' AS dest (TYPE sqlite)")
            con.execute("CREATE TABLE dest.data AS SELECT * FROM df_table")
        elif fmt == "docx":
            self._write_docx(df, path)
        elif fmt == "pdf":
            self._write_pdf(df, path)
        elif fmt in PPTX_FORMATS:
            self._write_pptx(df, path)
        elif fmt in ("jpg", "jpeg", "png"):
            # DataFrame → image: render as a styled table image
            self._write_df_as_image(df, path, fmt)
        else:
            df.write_csv(path)

    def write_excel_multi_sheet(self, sheets_dict: dict[str, pl.DataFrame], path: str) -> None:
        """Write multiple DataFrames as separate worksheets in one xlsx file."""
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        wb = openpyxl.Workbook()
        wb.remove(wb.active)                       # drop the default blank sheet
        HDR_FILL = PatternFill("solid", fgColor="4F46E5")
        HDR_FONT = Font(bold=True, color="FFFFFF")
        for sheet_name, df in sheets_dict.items():
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel sheet-name limit
            ws.append(list(df.columns))
            for cell in ws[1]:
                cell.fill = HDR_FILL
                cell.font = HDR_FONT
                cell.alignment = Alignment(horizontal="center")
            for row in df.iter_rows(named=False):
                ws.append([v if v is not None else "" for v in row])
        wb.save(path)

    def write_pdf_multi_sheet(self, sheets_dict: dict[str, pl.DataFrame], path: str) -> None:
        """Write multiple DataFrames as labelled sections in a single PDF."""
        import pymupdf as fitz
        import unicodedata
        import arabic_reshaper
        from bidi.algorithm import get_display

        FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
        FONT_BOLD    = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

        MARGIN   = 28
        ROW_H    = 16
        HDR_H    = 18
        SEC_H    = 20          # section-title bar height
        FONT_HDR = 8
        FONT_DAT = 7
        FONT_SEC = 10
        C_HDR_BG  = (79/255,  70/255,  229/255)
        C_HDR_TXT = (1.0, 1.0, 1.0)
        C_SEC_BG  = (30/255,  27/255,  75/255)
        C_EVEN    = (0.953, 0.957, 0.965)
        C_ODD     = (1.0,  1.0,  1.0)
        C_DAT_TXT = (0.118, 0.118, 0.118)
        C_BORDER  = (0.8,  0.8,  0.8)

        def _is_rtl(text: str) -> bool:
            rtl = sum(1 for ch in text if unicodedata.bidirectional(ch) in ("R", "AL", "AN"))
            ltr = sum(1 for ch in text if unicodedata.bidirectional(ch) == "L")
            return rtl > ltr

        def _prep(text: str) -> tuple[str, bool]:
            rtl = _is_rtl(text)
            return (get_display(arabic_reshaper.reshape(text)), True) if rtl else (text, False)

        # First pass: compute uniform page width based on max columns
        max_cols = max((len(df.columns) for df in sheets_dict.values()), default=1)
        PAGE_W, PAGE_H = (842, 595) if max_cols > 7 else (595, 842)
        rows_per_page  = int((PAGE_H - MARGIN * 2 - HDR_H) / ROW_H)

        doc  = fitz.open()

        def _new_page():
            pg = doc.new_page(width=PAGE_W, height=PAGE_H)
            pg.insert_font(fontname="dvr", fontfile=FONT_REGULAR)
            pg.insert_font(fontname="dvb", fontfile=FONT_BOLD)
            return pg

        def _draw_section(pg, y, title):
            col_w_full = PAGE_W - 2 * MARGIN
            rect = fitz.Rect(MARGIN, y, MARGIN + col_w_full, y + SEC_H)
            pg.draw_rect(rect, color=None, fill=C_SEC_BG, width=0, overlay=True)
            txt, rtl = _prep(title)
            pg.insert_textbox(
                fitz.Rect(rect.x0 + 4, rect.y0 + 3, rect.x1 - 4, rect.y1 - 3),
                txt, fontname="dvb", fontsize=FONT_SEC, color=C_HDR_TXT,
                align=fitz.TEXT_ALIGN_RIGHT if rtl else fitz.TEXT_ALIGN_LEFT,
                overlay=True,
            )
            return y + SEC_H + 4

        def _draw_row(pg, y, cells, col_w, n_cols, is_header, row_idx=0):
            h  = HDR_H if is_header else ROW_H
            fn = "dvb" if is_header else "dvr"
            fs = FONT_HDR if is_header else FONT_DAT
            tc = C_HDR_TXT if is_header else C_DAT_TXT
            for c_idx, text in enumerate(cells):
                x0   = MARGIN + c_idx * col_w
                rect = fitz.Rect(x0, y, x0 + col_w, y + h)
                fill = C_HDR_BG if is_header else (C_EVEN if row_idx % 2 == 0 else C_ODD)
                pg.draw_rect(rect, color=C_BORDER, fill=fill, width=0.4, overlay=True)
                inner = fitz.Rect(rect.x0 + 2, rect.y0 + 2, rect.x1 - 2, rect.y1 - 2)
                s, rtl = _prep(str(text)[:64])
                pg.insert_textbox(inner, s, fontname=fn, fontsize=fs, color=tc,
                                  align=fitz.TEXT_ALIGN_RIGHT if rtl else fitz.TEXT_ALIGN_LEFT,
                                  overlay=True)

        page = _new_page()
        y    = MARGIN

        for sheet_name, df in sheets_dict.items():
            n_cols    = len(df.columns)
            col_w     = max((PAGE_W - 2 * MARGIN) / n_cols, 40)
            n_rows    = min(len(df), 2000)
            row_data  = list(df.head(n_rows).iter_rows())
            col_names = list(df.columns)

            # If not enough room for section title + header + 1 data row → new page
            if y + SEC_H + HDR_H + ROW_H > PAGE_H - MARGIN:
                page = _new_page()
                y    = MARGIN

            y = _draw_section(page, y, sheet_name)
            _draw_row(page, y, col_names, col_w, n_cols, is_header=True)
            y += HDR_H

            data_idx = 0
            while data_idx < n_rows:
                if y + ROW_H > PAGE_H - MARGIN:
                    page = _new_page()
                    y    = MARGIN
                    _draw_row(page, y, col_names, col_w, n_cols, is_header=True)
                    y += HDR_H
                _draw_row(page, y,
                          [str(v) if v is not None else "" for v in row_data[data_idx]],
                          col_w, n_cols, is_header=False, row_idx=data_idx)
                y        += ROW_H
                data_idx += 1

            y += 10   # gap between sections

        doc.save(path)
        doc.close()

    def write_zip_multi_sheet(
        self, sheets_dict: dict[str, pl.DataFrame], base_path: str, fmt: str
    ) -> str:
        """Write each sheet as a separate file then zip them. Returns the zip path."""
        zip_path = str(base_path).rsplit(".", 1)[0] + ".zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for sheet_name, df in sheets_dict.items():
                safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in sheet_name)
                fname     = f"{safe_name}.{fmt}"
                with tempfile.NamedTemporaryFile(suffix=f".{fmt}", delete=False) as tmp:
                    tmp_path = tmp.name
                try:
                    self.write(df, tmp_path, fmt)
                    zf.write(tmp_path, fname)
                finally:
                    os.unlink(tmp_path)
        return zip_path

    def _write_docx(self, df: pl.DataFrame, path: str) -> None:
        from docx import Document
        doc = Document()
        table = doc.add_table(rows=1 + len(df), cols=len(df.columns))
        for i, col in enumerate(df.columns):
            table.rows[0].cells[i].text = col
        for r_idx, row in enumerate(df.iter_rows()):
            for c_idx, val in enumerate(row):
                table.rows[r_idx + 1].cells[c_idx].text = str(val) if val is not None else ""
        doc.save(path)

    def _write_pdf(self, df: pl.DataFrame, path: str) -> None:
        """Write DataFrame as a PDF table using PyMuPDF — full Unicode/Arabic support."""
        import pymupdf as fitz

        # DejaVu Sans ships with the Replit NixOS image and covers Arabic + Latin
        FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
        FONT_BOLD    = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

        n_cols = len(df.columns)
        n_rows = min(len(df), 2000)

        # Page size in points (A4): landscape for wide tables
        PAGE_W, PAGE_H = (842, 595) if n_cols > 7 else (595, 842)

        MARGIN   = 28
        ROW_H    = 16
        HDR_H    = 18
        FONT_HDR = 8
        FONT_DAT = 7
        col_w    = max((PAGE_W - 2 * MARGIN) / n_cols, 40)

        C_HDR_BG  = (79/255,  70/255,  229/255)
        C_HDR_TXT = (1.0, 1.0, 1.0)
        C_EVEN    = (0.953, 0.957, 0.965)
        C_ODD     = (1.0,  1.0,  1.0)
        C_DAT_TXT = (0.118, 0.118, 0.118)
        C_BORDER  = (0.8,  0.8,  0.8)

        rows_per_page = int((PAGE_H - MARGIN * 2 - HDR_H) / ROW_H)
        row_data  = list(df.head(n_rows).iter_rows())
        col_names = df.columns

        def _new_page(document: fitz.Document) -> fitz.Page:
            pg = document.new_page(width=PAGE_W, height=PAGE_H)
            pg.insert_font(fontname="dvr", fontfile=FONT_REGULAR)
            pg.insert_font(fontname="dvb", fontfile=FONT_BOLD)
            return pg

        import unicodedata
        import arabic_reshaper
        from bidi.algorithm import get_display

        def _is_rtl(text: str) -> bool:
            """Return True if the text is predominantly RTL (Arabic/Hebrew)."""
            rtl_count = 0
            ltr_count = 0
            for ch in text:
                bidi = unicodedata.bidirectional(ch)
                if bidi in ("R", "AL", "AN"):
                    rtl_count += 1
                elif bidi in ("L",):
                    ltr_count += 1
            return rtl_count > ltr_count

        def _prepare_text(text: str) -> tuple[str, bool]:
            """
            Return (display_text, is_rtl).
            For Arabic/RTL text: reshape letters and apply BiDi algorithm so the
            string renders correctly in a simple left-to-right PDF text stream.
            For LTR text: return as-is.
            """
            rtl = _is_rtl(text)
            if rtl:
                reshaped = arabic_reshaper.reshape(text)
                display  = get_display(reshaped)
            else:
                display = text
            return display, rtl

        def _draw_row(pg: fitz.Page, y: float, cells, is_header: bool, row_idx: int = 0) -> None:
            h          = HDR_H if is_header else ROW_H
            fn         = "dvb" if is_header else "dvr"
            fs         = FONT_HDR if is_header else FONT_DAT
            txt_color  = C_HDR_TXT if is_header else C_DAT_TXT
            for c_idx, text in enumerate(cells):
                x0         = MARGIN + c_idx * col_w
                rect       = fitz.Rect(x0, y, x0 + col_w, y + h)
                fill_color = C_HDR_BG if is_header else (C_EVEN if row_idx % 2 == 0 else C_ODD)
                pg.draw_rect(rect, color=C_BORDER, fill=fill_color, width=0.4, overlay=True)
                inner = fitz.Rect(rect.x0 + 2, rect.y0 + 2, rect.x1 - 2, rect.y1 - 2)
                cell_str            = str(text)[:64]
                display_str, is_rtl = _prepare_text(cell_str)
                align = fitz.TEXT_ALIGN_RIGHT if is_rtl else fitz.TEXT_ALIGN_LEFT
                pg.insert_textbox(
                    inner, display_str,
                    fontname=fn, fontsize=fs,
                    color=txt_color, align=align,
                    overlay=True,
                )

        doc      = fitz.open()
        page     = _new_page(doc)
        data_idx = 0

        while True:
            y = MARGIN
            _draw_row(page, y, col_names, is_header=True)
            y += HDR_H
            for _ in range(rows_per_page):
                if data_idx >= n_rows:
                    break
                _draw_row(page, y,
                          [str(v) if v is not None else "" for v in row_data[data_idx]],
                          is_header=False, row_idx=data_idx)
                y      += ROW_H
                data_idx += 1
            if data_idx >= n_rows:
                break
            page = _new_page(doc)

        doc.save(path)
        doc.close()

    def _write_pptx(self, df: pl.DataFrame, path: str) -> None:
        """Write DataFrame as a PPTX slide with a data table."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(slide_layout)

        rows_count = min(len(df) + 1, 50)  # header + data, cap at 50 rows per slide
        cols_count = len(df.columns)

        left   = Inches(0.5)
        top    = Inches(1.0)
        width  = Inches(9.0)
        height = Inches(0.4 * rows_count)

        table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

        # Header
        for c, col_name in enumerate(df.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x4F, 0x46, 0xE5)  # indigo
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Data rows
        for r_idx, row in enumerate(df.head(rows_count - 1).iter_rows()):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val) if val is not None else ""
                cell.text_frame.paragraphs[0].font.size = Pt(8)

        prs.save(path)

    def _write_df_as_image(self, df: pl.DataFrame, path: str, fmt: str) -> None:
        """Render a DataFrame as a simple table image using PIL."""
        from PIL import Image, ImageDraw, ImageFont

        cell_w, cell_h = 120, 28
        n_cols = len(df.columns)
        n_rows = min(len(df), 50) + 1  # +1 for header

        img_w = cell_w * n_cols + 2
        img_h = cell_h * n_rows + 2

        img = Image.new("RGB", (img_w, img_h), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.load_default(size=12)
        except TypeError:
            font = ImageFont.load_default()

        # Header
        for c, col in enumerate(df.columns):
            x, y = c * cell_w + 1, 1
            draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], fill=(79, 70, 229))
            draw.text((x + 4, y + 6), str(col)[:16], fill=(255, 255, 255), font=font)

        # Data rows
        for r_idx, row in enumerate(df.head(50).iter_rows()):
            bg = (249, 250, 251) if r_idx % 2 == 0 else (255, 255, 255)
            for c_idx, val in enumerate(row):
                x = c_idx * cell_w + 1
                y = (r_idx + 1) * cell_h + 1
                draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], fill=bg)
                draw.text((x + 4, y + 6), str(val if val is not None else "")[:16],
                          fill=(30, 30, 30), font=font)

        save_fmt = "JPEG" if fmt in ("jpg", "jpeg") else "PNG"
        img.save(path, save_fmt)

    # ─── Direct (non-tabular) conversions ────────────────────────────────────

    def convert_direct(self, src_path: str, src_fmt: str, dst_path: str, dst_fmt: str) -> str:
        """
        Convert between formats that don't pass through a DataFrame.
        Returns the actual output path (may differ if multi-page → zip).
        """
        src_fmt = src_fmt.lower().lstrip(".")
        dst_fmt = dst_fmt.lower().lstrip(".")

        # Image → PDF
        if src_fmt in IMAGE_FORMATS and dst_fmt == "pdf":
            self._image_to_pdf(src_path, dst_path)
            return dst_path

        # PDF → image(s)
        if src_fmt == "pdf" and dst_fmt in ("jpg", "jpeg", "png"):
            return self._pdf_to_images(src_path, dst_path, dst_fmt)

        # SVG → PDF
        if src_fmt == "svg" and dst_fmt == "pdf":
            self._svg_to_pdf(src_path, dst_path)
            return dst_path

        # PDF → SVG (first page)
        if src_fmt == "pdf" and dst_fmt == "svg":
            self._pdf_to_svg(src_path, dst_path)
            return dst_path

        raise ValueError(f"No direct conversion path for {src_fmt} → {dst_fmt}")

    def _image_to_pdf(self, src_path: str, dst_path: str) -> None:
        """Embed an image into an A4 PDF page, centred and scaled to fit."""
        from PIL import Image
        from fpdf import FPDF

        img = Image.open(src_path)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        # Save as temporary JPEG for fpdf
        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            img.save(tmp_path, "JPEG", quality=92)

            dpi = 96
            w_mm = img.size[0] / dpi * 25.4
            h_mm = img.size[1] / dpi * 25.4

            # Fit inside A4 margins
            max_w, max_h = 190.0, 270.0
            if w_mm > max_w or h_mm > max_h:
                scale = min(max_w / w_mm, max_h / h_mm)
                w_mm, h_mm = w_mm * scale, h_mm * scale

            pdf = FPDF(unit="mm", format="A4")
            pdf.add_page()
            x = (210.0 - w_mm) / 2
            y = (297.0 - h_mm) / 2
            pdf.image(tmp_path, x=x, y=y, w=w_mm, h=h_mm)
            pdf.output(dst_path)
        finally:
            os.unlink(tmp_path)

    def _pdf_to_images(self, src_path: str, dst_path: str, fmt: str) -> str:
        """
        Render each PDF page as an image.
        - Single page  → writes image directly to dst_path.
        - Multi-page   → writes a ZIP archive (dst_path renamed to .zip).
        Returns the actual path written.
        """
        import fitz  # pymupdf

        save_fmt = "jpeg" if fmt in ("jpg", "jpeg") else "png"
        doc = fitz.open(src_path)
        mat = fitz.Matrix(2.0, 2.0)  # ~144 DPI

        if len(doc) == 1:
            pix = doc[0].get_pixmap(matrix=mat, alpha=False)
            pix.save(dst_path)
            return dst_path

        # Multi-page: produce ZIP
        zip_path = str(Path(dst_path).with_suffix(".zip"))
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_bytes = pix.tobytes(save_fmt)
                zf.writestr(f"page_{i + 1:03d}.{fmt}", img_bytes)
        return zip_path

    def _svg_to_pdf(self, src_path: str, dst_path: str) -> None:
        """Convert SVG to PDF using PyMuPDF."""
        import fitz
        doc = fitz.open(src_path)   # fitz can open SVG directly
        pdf_bytes = doc.convert_to_pdf()
        Path(dst_path).write_bytes(pdf_bytes)

    def _pdf_to_svg(self, src_path: str, dst_path: str) -> None:
        """Export the first page of a PDF as an SVG file."""
        import fitz
        doc = fitz.open(src_path)
        page = doc[0]
        svg_text = page.get_svg_image()
        Path(dst_path).write_text(svg_text, encoding="utf-8")

    # ─── Preview ──────────────────────────────────────────────────────────────

    def preview(self, path: str, fmt: str, rows: int = 100) -> dict:
        """Return preview data: columns, sample rows, shape."""
        try:
            df = self.read(path, fmt)
            total_rows, total_cols = df.shape
            sample = df.head(rows)
            return {
                "columns": df.columns,
                "dtypes": [str(d) for d in df.dtypes],
                "rows": sample.to_dicts(),
                "total_rows": total_rows,
                "total_cols": total_cols,
                "shape": f"{total_rows:,} × {total_cols}",
            }
        except Exception as e:
            logger.error(f"Preview failed for {path}: {e}")
            return {"error": str(e), "columns": [], "rows": [], "total_rows": 0, "total_cols": 0}

    # ─── Metadata ─────────────────────────────────────────────────────────────

    def get_metadata(self, path: str, fmt: str) -> dict:
        """Extract lightweight metadata without reading the full file."""
        meta: dict[str, Any] = {}
        try:
            if fmt in EXCEL_FORMATS:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                meta["sheets"] = wb.sheetnames
                ws = wb.active
                meta["rows"] = ws.max_row
                meta["columns"] = ws.max_column
                wb.close()
            elif fmt in CSV_LIKE:
                sep = "\t" if fmt == "tsv" else ","
                df = pl.scan_csv(path, separator=sep, infer_schema_length=100).limit(1).collect()
                meta["columns"] = len(df.columns)
                with open(path, "rb") as f:
                    meta["rows"] = sum(1 for _ in f)
            elif fmt == "parquet":
                import pyarrow.parquet as pq
                pf = pq.ParquetFile(path)
                meta["rows"] = pf.metadata.num_rows
                meta["columns"] = pf.metadata.num_columns
            elif fmt in IMAGE_FORMATS:
                from PIL import Image
                img = Image.open(path)
                meta["width"], meta["height"] = img.size
                meta["mode"] = img.mode
            elif fmt == "pdf":
                import fitz
                doc = fitz.open(path)
                meta["pages"] = len(doc)
        except Exception as e:
            logger.debug(f"Partial metadata extraction: {e}")
        return meta

    # ─── Sheet listing ────────────────────────────────────────────────────────

    def list_sheets(self, path: str, fmt: str) -> list[str]:
        if fmt not in EXCEL_FORMATS:
            return []
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True)
            sheets = wb.sheetnames
            wb.close()
            return sheets
        except Exception:
            return []
